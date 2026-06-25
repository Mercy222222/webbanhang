from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def set_slide_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(10, 15, 28) # Dark Navy/Black background

def add_title(slide, text, color_pink):
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = color_pink

def create_content_textbox(slide):
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf

def create_thesis_presentation():
    prs = Presentation()
    
    # Define Theme Colors
    color_pink = RGBColor(236, 72, 153)
    color_blue = RGBColor(59, 130, 246)
    color_white = RGBColor(255, 255, 255)
    color_gray = RGBColor(170, 170, 170)
    color_green = RGBColor(16, 185, 129)

    # ----------------------------------------------------
    # 1. TITLE SLIDE
    # ----------------------------------------------------
    slide_layout = prs.slide_layouts[6] # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = "BÁO CÁO ĐỒ ÁN"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = color_gray
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "WEBSITE BÁN BÁNH NGỌT"
    p2.font.size = Pt(50)
    p2.font.bold = True
    p2.font.color.rgb = color_pink
    p2.alignment = PP_ALIGN.CENTER
    
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(2))
    tf2 = txBox2.text_frame
    p3 = tf2.add_paragraph()
    p3.text = "Môn: Phân Tích Thiết Kế Hệ Thống\nGVHD: Th.S Nguyễn Văn Danh\n"
    p3.font.size = Pt(20)
    p3.font.color.rgb = color_blue
    p3.alignment = PP_ALIGN.CENTER
    
    p4 = tf2.add_paragraph()
    p4.text = "Sinh viên: Hoàng Minh Triều | Nguyễn Gia Sinh\nVõ Hữu Trí | Phan Tấn Phát"
    p4.font.size = Pt(16)
    p4.font.color.rgb = color_white
    p4.alignment = PP_ALIGN.CENTER

    # ----------------------------------------------------
    # 2. OVERVIEW SLIDE
    # ----------------------------------------------------
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    add_title(slide, "CHƯƠNG 1: TỔNG QUAN ĐỀ TÀI", color_pink)
    
    tf = create_content_textbox(slide)
    
    p = tf.add_paragraph()
    p.text = "Lý do chọn đề tài:"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = color_blue
    
    p = tf.add_paragraph()
    p.text = "• Hành vi người dùng chuyển dịch mạnh sang mua sắm trực tuyến.\n• Khắc phục hạn chế của cửa hàng bán bánh truyền thống.\n• Tối ưu hóa chu trình: Nhận đơn - Làm bánh - Giao hàng."
    p.font.size = Pt(20)
    p.font.color.rgb = color_white
    
    p = tf.add_paragraph()
    p.text = "\nMục tiêu hệ thống:"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = color_blue
    
    p = tf.add_paragraph()
    p.text = "• Hệ thống kinh doanh trực tuyến: Giao diện trực quan, bộ lọc thông minh.\n• Hệ thống vận hành nội bộ (Kho, Bếp, Bán hàng) minh bạch.\n• Frontend: HTML/CSS/JS (React/Vue). Backend: NodeJS/PHP. DB: MySQL."
    p.font.size = Pt(20)
    p.font.color.rgb = color_white

    # ----------------------------------------------------
    # 3. 6 BIỂU MẪU (FORMS)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    add_title(slide, "CHƯƠNG 2: HỆ THỐNG BIỂU MẪU NGHIỆP VỤ", color_pink)
    
    tf = create_content_textbox(slide)
    
    p = tf.add_paragraph()
    p.text = "Dựa trên khảo sát thực tế (Tiệm Bánh Làng Xì Mi), hệ thống áp dụng 6 biểu mẫu lõi:"
    p.font.size = Pt(22)
    p.font.color.rgb = color_gray
    p.font.italic = True
    
    forms = [
        "1. Biểu mẫu Bảng thanh toán tiền lương (Mẫu 02-L)",
        "2. Biểu mẫu Phiếu thu (Mẫu 01-TT)",
        "3. Biểu mẫu Phiếu nhập kho (Mẫu 01-VT)",
        "4. Biểu mẫu Chứng từ thanh toán lương (Mẫu CT-TTL)",
        "5. Biểu mẫu Phiếu chi (Mẫu 02-TT)",
        "6. Biểu mẫu Phiếu xuất kho (Mẫu 02-VT)"
    ]
    
    for idx, form in enumerate(forms):
        p = tf.add_paragraph()
        p.text = f"\n{form}"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = color_white if idx % 2 == 0 else color_blue

    # ----------------------------------------------------
    # 4. SƠ ĐỒ BFD & USE CASE
    # ----------------------------------------------------
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    add_title(slide, "SƠ ĐỒ CHỨC NĂNG (BFD) & USE CASE", color_pink)
    
    tf = create_content_textbox(slide)
    
    p = tf.add_paragraph()
    p.text = "1. Sơ đồ BFD (Business Function Diagram)"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = color_blue
    
    p = tf.add_paragraph()
    p.text = "Phân rã các chức năng thành các khối nghiệp vụ logic độc lập."
    p.font.size = Pt(18)
    p.font.color.rgb = color_gray
    
    p = tf.add_paragraph()
    p.text = "\n2. Sơ đồ Use Case (Mô hình hóa tác nhân và hành vi)"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = color_blue
    
    p = tf.add_paragraph()
    p.text = "• Sơ đồ Tổng quát\n• Sơ đồ Phân rã Cấp 1 (8 modules):"
    p.font.size = Pt(20)
    p.font.color.rgb = color_white
    
    p = tf.add_paragraph()
    p.text = "   - Quản lý sản phẩm (Bánh)  |  - Quản lý đơn hàng\n   - Quản lý tài khoản        |  - Quản lý giỏ hàng\n   - Phân rã tìm kiếm         |  - Quản lý nhân viên\n   - Phân rã xem sản phẩm     |  - Quản lý người dùng"
    p.font.size = Pt(18)
    p.font.color.rgb = color_green
    
    p = tf.add_paragraph()
    p.text = "• Sơ đồ Phân rã Cấp 2 (4 quy trình lõi):"
    p.font.size = Pt(20)
    p.font.color.rgb = color_white
    
    p = tf.add_paragraph()
    p.text = "   - Phân rã Đăng nhập        |  - Phân rã Đặt hàng\n   - Phân rã Thêm bánh        |  - Phân rã Trả hàng"
    p.font.size = Pt(18)
    p.font.color.rgb = color_green

    # ----------------------------------------------------
    # 5. SƠ ĐỒ ACTIVITY & SEQUENCE
    # ----------------------------------------------------
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    add_title(slide, "SƠ ĐỒ HOẠT ĐỘNG (ACTIVITY) & TUẦN TỰ (SEQUENCE)", color_pink)
    
    tf = create_content_textbox(slide)
    
    p = tf.add_paragraph()
    p.text = "Thiết kế chi tiết luồng xử lý và tương tác hệ thống qua 4 luồng chính:"
    p.font.size = Pt(22)
    p.font.color.rgb = color_white
    
    flows = [
        "Luồng Đăng nhập (Xác thực và phân quyền)",
        "Luồng Đặt hàng (Quy trình Checkout, giỏ hàng, xác nhận)",
        "Luồng Trả hàng (Xử lý hoàn trả, cập nhật kho)",
        "Luồng Thêm bánh (Nhập liệu từ Admin, lưu vào Database)"
    ]
    
    for flow in flows:
        p = tf.add_paragraph()
        p.text = f"\n✔️ {flow}"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = color_blue
        
        p = tf.add_paragraph()
        p.text = "   - Cung cấp sơ đồ Hoạt Động (Activity Diagram)\n   - Cung cấp sơ đồ Tuần Tự (Sequence Diagram)"
        p.font.size = Pt(18)
        p.font.color.rgb = color_gray

    # ----------------------------------------------------
    # 6. ENDING SLIDE
    # ----------------------------------------------------
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    
    txBox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.add_paragraph()
    p.text = "XIN CẢM ƠN!"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = color_pink
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "Phần trình bày báo cáo Website Bán Bánh Ngọt kết thúc."
    p2.font.size = Pt(24)
    p2.font.color.rgb = color_white
    p2.alignment = PP_ALIGN.CENTER
    
    prs.save('Bao_Cao_Website_Ban_Banh.pptx')
    print("Presentation revised and saved successfully!")

if __name__ == '__main__':
    create_thesis_presentation()
