from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    
    # 1. Title Slide
    slide_layout = prs.slide_layouts[0] # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "MERCY SMM PANEL"
    title.text_frame.paragraphs[0].font.size = Pt(54)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(236, 72, 153) # Pink
    
    subtitle.text = "The Best & Cheapest SMM Panel for Resellers\nGrow Your Social Media Accounts With Most Trusted Panel"
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)

    # 2. About Us Slide
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "About Mercy SMM Panel"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(59, 130, 246) # Blue
    
    tf = content.text_frame
    tf.text = "Welcome to the most trusted provider in the market."
    p = tf.add_paragraph()
    p.text = "• We provide the fastest social media marketing services."
    p = tf.add_paragraph()
    p.text = "• High-quality real engagement for your social accounts."
    p = tf.add_paragraph()
    p.text = "• Automated system processing thousands of orders per minute."
    p = tf.add_paragraph()
    p.text = "• Secure, anonymous, and reseller-friendly."

    # 3. Our Services
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Our Premium Services"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(139, 92, 246) # Purple
    
    tf = content.text_frame
    tf.text = "We cover all major platforms:"
    tf.add_paragraph().text = "• Facebook: Likes, Followers, Views, Comments"
    tf.add_paragraph().text = "• Instagram: Followers, Story Views, Reel Views"
    tf.add_paragraph().text = "• TikTok: Views, Followers, Likes, Shares"
    tf.add_paragraph().text = "• YouTube: Subscribers, Watch Time, Views"
    tf.add_paragraph().text = "• Twitter/X & Telegram Services"

    # 4. Statistics
    slide_layout = prs.slide_layouts[5] # Title only
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Platform Statistics"
    
    # Add custom text boxes for stats
    left = Inches(1)
    top = Inches(3)
    width = Inches(8)
    height = Inches(2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.add_paragraph()
    p.text = "1,402,154+ \nTotal Orders Completed\n\n"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(16, 185, 129) # Green
    
    p2 = tf.add_paragraph()
    p2.text = "1,514 Active Services    |    50,000+ Happy Users"
    p2.font.size = Pt(28)
    p2.font.bold = True

    # 5. Call To Action
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Ready to Grow?"
    title.text_frame.paragraphs[0].font.size = Pt(50)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(236, 72, 153)
    
    subtitle.text = "Register for free today.\nwww.mercysmmpanel.com"
    subtitle.text_frame.paragraphs[0].font.size = Pt(30)
    
    prs.save('Mercy_SMM_Panel_Pitch_Deck.pptx')
    print("Presentation saved successfully as Mercy_SMM_Panel_Pitch_Deck.pptx")

if __name__ == '__main__':
    create_presentation()
