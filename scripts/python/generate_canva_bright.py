from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def set_slide_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(244, 247, 254) # Light grayish blue

def add_title(slide, text, color_title):
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = color_title

def create_content_textbox(slide, top=1.8):
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(9), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf

def create_image_placeholder(slide, text="Kéo thả hình ảnh Sơ đồ/Biểu mẫu vào đây trên Canva"):
    # Create a rectangle shape to act as an image placeholder for Canva
    left = Inches(5.5)
    top = Inches(1.8)
    width = Inches(4)
    height = Inches(4.5)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255) # White bg
    shape.line.color.rgb = RGBColor(79, 70, 229) # Indigo border
    shape.line.width = Pt(2)
    
    tf = shape.text_frame
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(148, 163, 184) # Slate 400
    p.alignment = PP_ALIGN.CENTER

def create_bright_canva_presentation():
    prs = Presentation()
    
    # Theme Colors matching the Bright Corporate template
    color_indigo = RGBColor(67, 56, 202) # Title / Primary
    color_slate = RGBColor(30, 41, 59) # Body text
    color_blue = RGBColor(59, 130, 246)
    color_pink = RGBColor(236, 72, 153)
    color_white = RGBColor(255, 255, 255)

    # ----------------------------------------------------
    # 1. TITLE SLIDE
    # ----------------------------------------------------
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = "BÁO CÁO ĐỒ ÁN MÔN HỌC"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = color_slate
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "WEBSITE BÁN BÁNH NGỌT"
    p2.font.size = Pt(54)
    p2.font.bold = True
    p2.font.color.rgb = color_indigo
    p2.alignment = PP_ALIGN.CENTER
    
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(2))
    tf2 = txBox2.text_frame
    p3 = tf2.add_paragraph()
    p3.text = "Môn: Phân Tích Thiết Kế Hệ Thống\nGVHD: Th.S Nguyễn Văn Danh\nNhóm SV: Hoàng Minh Triều, Nguyễn Gia Sinh, Võ Hữu Trí, Phan Tấn Phát"
    p3.font.size = Pt(20)
    p3.font.color.rgb = color_slate
    p3.alignment = PP_ALIGN.CENTER

    # ----------------------------------------------------
    # 2. OVERVIEW SLIDE
    # ----------------------------------------------------
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    add_title(slide, "1. TỔNG QUAN ĐỀ TÀI", color_indigo)
    tf = create_content_textbox(slide)
    
    p = tf.add_paragraph()
    p.text = "Lý do chọn đề tài:"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = color_blue
    
    p = tf.add_paragraph()
    p.text = "• Thay đổi hành vi mua sắm sang trực tuyến.\n• Khắc phục giới hạn của mô hình truyền thống.\n• Tối ưu hóa chu trình: Đặt hàng - Bếp - Giao nhận khép kín."
    p.font.size = Pt(20)
    p.font.color.rgb = color_slate
    
    p = tf.add_paragraph()
    p.text = "\nMục tiêu hệ thống:"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = color_pink
    
    p = tf.add_paragraph()
    p.text = "• Khách hàng: Đặt bánh nhanh, tùy chỉnh lời chúc, theo dõi đơn.\n• Quản lý: Giám sát kho, định mức nguyên liệu, POS, báo cáo tự động."
    p.font.size = Pt(20)
    p.font.color.rgb = color_slate

    # ----------------------------------------------------
    # 3. 6 BIỂU MẪU (FORMS)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    add_title(slide, "2. HỆ THỐNG 6 BIỂU MẪU NGHIỆP VỤ", color_indigo)
    tf = create_content_textbox(slide)
    
    p = tf.add_paragraph()
    p.text = "Các biểu mẫu lõi được số hóa từ quy trình thực tế:"
    p.font.size = Pt(22)
    p.font.color.rgb = color_slate
    p.font.italic = True
    
    forms = [
        "1. Bảng thanh toán tiền lương (Mẫu 02-L)",
        "2. Phiếu thu (Mẫu 01-TT)",
        "3. Phiếu nhập kho (Mẫu 01-VT)",
        "4. Chứng từ thanh toán lương (Mẫu CT-TTL)",
        "5. Phiếu chi (Mẫu 02-TT)",
        "6. Phiếu xuất kho (Mẫu 02-VT)"
    ]
    
    for form in forms:
        p = tf.add_paragraph()
        p.text = f"✔️ {form}"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = color_indigo
        
    create_image_placeholder(slide, "Canva: Kéo thả ảnh Biểu mẫu vào đây")

    # ----------------------------------------------------
    # 4. BFD DIAGRAM
    # ----------------------------------------------------
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    add_title(slide, "3. SƠ ĐỒ CHỨC NĂNG (BFD)", color_indigo)
    tf = create_content_textbox(slide)
    
    p = tf.add_paragraph()
    p.text = "Business Function Diagram"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = color_blue
    
    p = tf.add_paragraph()
    p.text = "\nHệ thống được phân rã thành các phân hệ nghiệp vụ độc lập:"
    p.font.size = Pt(20)
    p.font.color.rgb = color_slate
    
    p = tf.add_paragraph()
    p.text = "• Quản lý Bán hàng\n• Quản lý Kho & Nguyên liệu\n• Quản lý Nhân sự\n• Quản trị Hệ thống (Admin)"
    p.font.size = Pt(20)
    p.font.color.rgb = color_indigo
    
    create_image_placeholder(slide, "Canva: Kéo thả ảnh Sơ đồ BFD vào đây")

    # ----------------------------------------------------
    # 5. USE CASE - CẤP 1
    # ----------------------------------------------------
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    add_title(slide, "4. SƠ ĐỒ USE CASE (TỔNG QUÁT & CẤP 1)", color_indigo)
    tf = create_content_textbox(slide)
    
    p = tf.add_paragraph()
    p.text = "Phân rã Cấp 1 (8 Khối chức năng lõi):"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = color_blue
    
    usecases1 = [
        "Quản lý Sản phẩm (Bánh)", "Quản lý Đơn hàng",
        "Quản lý Tài khoản", "Quản lý Giỏ hàng",
        "Tìm kiếm", "Xem sản phẩm",
        "Quản lý Nhân viên", "Quản lý Người dùng"
    ]
    
    for uc in usecases1:
        p = tf.add_paragraph()
        p.text = f"• {uc}"
        p.font.size = Pt(20)
        p.font.color.rgb = color_slate
        
    create_image_placeholder(slide, "Canva: Kéo thả ảnh Sơ đồ Use Case Cấp 1")

    # ----------------------------------------------------
    # 6. USE CASE - CẤP 2
    # ----------------------------------------------------
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    add_title(slide, "5. SƠ ĐỒ USE CASE (PHÂN RÃ CẤP 2)", color_indigo)
    tf = create_content_textbox(slide)
    
    p = tf.add_paragraph()
    p.text = "Phân rã Cấp 2 (4 Quy trình chi tiết):"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = color_pink
    
    usecases2 = [
        "Đăng nhập (Xác thực, phân quyền)",
        "Đặt hàng (Checkout, cập nhật giỏ)",
        "Thêm bánh (Upload, cập nhật giá, tồn kho)",
        "Trả hàng (Xử lý hoàn trả, cộng lại kho)"
    ]
    
    for uc in usecases2:
        p = tf.add_paragraph()
        p.text = f"• {uc}"
        p.font.size = Pt(22)
        p.font.color.rgb = color_slate
        p.space_after = Pt(10)
        
    create_image_placeholder(slide, "Canva: Kéo thả ảnh Sơ đồ Use Case Cấp 2")

    # ----------------------------------------------------
    # 7. ACTIVITY DIAGRAM
    # ----------------------------------------------------
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    add_title(slide, "6. SƠ ĐỒ HOẠT ĐỘNG (ACTIVITY DIAGRAM)", color_indigo)
    tf = create_content_textbox(slide)
    
    p = tf.add_paragraph()
    p.text = "Mô hình hóa luồng xử lý hệ thống:"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = color_blue
    
    acts = [
        "Luồng Đăng nhập",
        "Luồng Đặt hàng",
        "Luồng Trả hàng",
        "Luồng Thêm bánh"
    ]
    
    for act in acts:
        p = tf.add_paragraph()
        p.text = f"🔄 {act}"
        p.font.size = Pt(24)
        p.font.color.rgb = color_slate
        p.space_after = Pt(5)
        
    create_image_placeholder(slide, "Canva: Kéo thả ảnh Sơ đồ Activity vào đây")

    # ----------------------------------------------------
    # 8. SEQUENCE DIAGRAM
    # ----------------------------------------------------
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    add_title(slide, "7. SƠ ĐỒ TUẦN TỰ (SEQUENCE DIAGRAM)", color_indigo)
    tf = create_content_textbox(slide)
    
    p = tf.add_paragraph()
    p.text = "Mô hình hóa tương tác giữa các đối tượng (Actors & System):"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = color_pink
    
    seqs = [
        "Sequence Đăng nhập",
        "Sequence Đặt hàng",
        "Sequence Trả hàng",
        "Sequence Thêm bánh"
    ]
    
    for seq in seqs:
        p = tf.add_paragraph()
        p.text = f"⇄ {seq}"
        p.font.size = Pt(24)
        p.font.color.rgb = color_slate
        p.space_after = Pt(5)
        
    create_image_placeholder(slide, "Canva: Kéo thả ảnh Sơ đồ Sequence vào đây")

    # ----------------------------------------------------
    # 9. ENDING SLIDE
    # ----------------------------------------------------
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = "XIN CẢM ƠN!"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = color_indigo
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "Thầy và các bạn đã chú ý lắng nghe."
    p2.font.size = Pt(24)
    p2.font.color.rgb = color_slate
    p2.alignment = PP_ALIGN.CENTER
    
    prs.save('Bao_Cao_Website_Ban_Banh_Canva_Bright.pptx')
    print("Canva Bright PPT saved successfully!")

if __name__ == '__main__':
    create_bright_canva_presentation()
